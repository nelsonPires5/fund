#!/usr/bin/env python3
from pathlib import Path
import json, math
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import MSO_ANCHOR

RUN = Path('sectors/software/companies/SHOP/runs/2026-05-23-initial-coverage')
CHART_DIR = RUN / 'assets' / 'charts'
CHART_DIR.mkdir(parents=True, exist_ok=True)
OUT = RUN / 'deck.pptx'
outputs = json.loads((RUN/'outputs.json').read_text())
out = {x['key']: x for x in outputs['outputs']}

# Data from report/model/Koyfin summary
price = out['market.current_price']['value']
dcf_value = out['valuation.implied_share_price']['value']
upside = out['valuation.upside_downside']['value']
enterprise_value = out['valuation.enterprise_value']['value'] / 1000
wacc = out['assumption.wacc']['value']
term_g = out['assumption.terminal_growth']['value']

hist_years = [2022, 2023, 2024, 2025]
hist_rev = [5.6, 7.06, 8.88, 11.56]
hist_fcf = [-0.19, 0.91, 1.60, 2.01]
proj_years = [2026, 2027, 2028, 2029, 2030, 2031, 2032, 2033, 2034, 2035]
proj_rev = [14.80,18.32,22.93,27.52,32.20,36.71,40.75,44.42,47.53,49.91]
proj_fcf = [2.44,3.21,4.24,5.23,6.28,7.34,8.35,9.33,10.22,10.98]

# Chart styling
NAVY = '#1F4E79'
BLUE = '#2F75B5'
LTBLUE = '#BDD7EE'
GREY = '#666666'
LTGREY = '#F2F2F2'
GREEN = '#63BE7B'
RED = '#C00000'
ORANGE = '#F4B183'

plt.rcParams['font.family'] = 'DejaVu Sans'
plt.rcParams['axes.edgecolor'] = '#BBBBBB'
plt.rcParams['axes.labelcolor'] = '#333333'
plt.rcParams['xtick.color'] = '#333333'
plt.rcParams['ytick.color'] = '#333333'

# 1) Revenue/FCF chart
fig, ax1 = plt.subplots(figsize=(9,4.8))
years = hist_years + proj_years
rev = hist_rev + proj_rev
fcf = hist_fcf + proj_fcf
colors = ['#9EADCC']*len(hist_years) + [BLUE]*len(proj_years)
ax1.bar([str(y) for y in years], rev, color=colors, label='Revenue')
ax1.set_ylabel('Revenue ($B)')
ax1.tick_params(axis='x', rotation=45)
ax2 = ax1.twinx()
ax2.plot([str(y) for y in years], fcf, color=GREEN, marker='o', linewidth=2.5, label='FCF')
ax2.set_ylabel('FCF ($B)')
ax1.set_title('Revenue scales from $11.6B in FY25 to ~$49.9B in FY35E; FCF expands with operating leverage')
ax1.grid(axis='y', alpha=0.25)
fig.tight_layout()
rev_chart = CHART_DIR/'shop_revenue_fcf_projection.png'
fig.savefig(rev_chart, dpi=180, bbox_inches='tight')
plt.close(fig)

# 2) Valuation range chart
fig, ax = plt.subplots(figsize=(9,3.6))
ranges=[('Bear',60,85),('Base DCF',90,105),('Upside base',110,135),('Bull',150,200)]
ypos=range(len(ranges))
for i,(name,lo,hi) in enumerate(ranges):
    ax.barh(i, hi-lo, left=lo, color=[RED, BLUE, LTBLUE, GREEN][i], alpha=0.85)
    ax.text((lo+hi)/2, i, f'${lo}–${hi}', va='center', ha='center', color='white' if i in [0,1,3] else '#222', weight='bold')
ax.axvline(price, color='black', linestyle='--', linewidth=1.5, label=f'Current ${price:.0f}')
ax.axvline(dcf_value, color=NAVY, linestyle='-', linewidth=1.8, label=f'Model base ${dcf_value:.0f}')
ax.set_yticks(list(ypos)); ax.set_yticklabels([x[0] for x in ranges])
ax.set_xlabel('Value per share')
ax.set_xlim(35,210)
ax.set_title('Valuation range: quality is high, but margin of safety is limited at the current price')
ax.legend(loc='lower right', frameon=False)
ax.grid(axis='x', alpha=0.25)
fig.tight_layout()
val_chart = CHART_DIR/'shop_valuation_range.png'
fig.savefig(val_chart, dpi=180, bbox_inches='tight')
plt.close(fig)

# 3) KPI scorecard chart
fig, ax = plt.subplots(figsize=(9,4))
kpis=['Revenue\ngrowth','GMV\ngrowth','GPV\npenetration','FCF\nmargin','Gross\nmargin','EV/Sales\nNTM']
vals=[34,35,67,15,48,8.3]
labels=['+34%','+35%','67%','15%','48%','8.3x']
bar_colors=[GREEN,GREEN,BLUE,BLUE,BLUE,ORANGE]
ax.bar(kpis, vals, color=bar_colors)
for i,v in enumerate(vals):
    ax.text(i, v+max(vals)*0.025, labels[i], ha='center', fontweight='bold')
ax.set_ylim(0,80)
ax.set_title('KPI dashboard: fundamentals strong; valuation still demanding')
ax.set_ylabel('% / multiple')
ax.grid(axis='y', alpha=0.2)
fig.tight_layout()
kpi_chart = CHART_DIR/'shop_kpi_dashboard.png'
fig.savefig(kpi_chart, dpi=180, bbox_inches='tight')
plt.close(fig)

# 4) Multiple sanity check chart
fig, ax = plt.subplots(figsize=(8.5,4.2))
labels=['SHOP\nLTM','SHOP\nNTM','Wix\nNTM','E-comm SW\nmedian','SaaS\nmedian']
vals=[10.3,8.3,2.4,1.3,2.5]
ax.bar(labels, vals, color=[NAVY,BLUE,'#7F9CCB','#BFBFBF','#BFBFBF'])
for i,v in enumerate(vals): ax.text(i,v+0.25,f'{v:.1f}x',ha='center',fontweight='bold')
ax.set_ylim(0,12)
ax.set_ylabel('EV / Revenue')
ax.set_title('SHOP still trades at a large premium even after the pullback')
ax.grid(axis='y', alpha=0.2)
fig.tight_layout()
mult_chart = CHART_DIR/'shop_ev_sales_comps.png'
fig.savefig(mult_chart, dpi=180, bbox_inches='tight')
plt.close(fig)

# 5) WACC/g heatmap using the updated model's mid-year convention and base FCF forecast
wacc_vals=[8.5,9.5,10.5]
g_vals=[2.5,3.5,4.5]
def model_price_for(wacc_pct, g_pct):
    w = wacc_pct/100.0; g = g_pct/100.0
    pv = sum(f/((1+w)**(0.5+i)) for i, f in enumerate(proj_fcf))
    pvtv = (proj_fcf[-1]*(1+g)/(w-g))/((1+w)**9.5)
    return (pv+pvtv+5.561)/1.3
heat=[[model_price_for(w,g) for g in g_vals] for w in wacc_vals]
fig, ax = plt.subplots(figsize=(5.5,3.8))
im=ax.imshow(heat, cmap='RdYlGn', aspect='auto')
ax.set_xticks(range(3)); ax.set_xticklabels([f'{g:.1f}%' for g in g_vals])
ax.set_yticks(range(3)); ax.set_yticklabels([f'{w:.1f}%' for w in wacc_vals])
ax.set_xlabel('Terminal growth'); ax.set_ylabel('WACC')
for i in range(3):
    for j in range(3): ax.text(j,i,f'${heat[i][j]:.0f}',ha='center',va='center',fontweight='bold')
ax.set_title('DCF sensitivity: discount rate and terminal assumptions drive value')
fig.colorbar(im, ax=ax, fraction=0.046, pad=0.04)
fig.tight_layout()
heat_chart = CHART_DIR/'shop_wacc_g_sensitivity.png'
fig.savefig(heat_chart, dpi=180, bbox_inches='tight')
plt.close(fig)

# PPT helpers
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

RGB_NAVY = RGBColor(31,78,121); RGB_BLUE=RGBColor(47,117,181); RGB_LTGREY=RGBColor(242,242,242); RGB_WHITE=RGBColor(255,255,255)
RGB_DARK=RGBColor(40,40,40); RGB_GREEN=RGBColor(0,128,0); RGB_RED=RGBColor(192,0,0)

def add_textbox(slide, x, y, w, h, text, font_size=16, bold=False, color=RGB_DARK, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear(); tf.word_wrap=True; tf.margin_left=Pt(3); tf.margin_right=Pt(3); tf.margin_top=Pt(2); tf.margin_bottom=Pt(2)
    p = tf.paragraphs[0]
    p.text = text
    if align: p.alignment=align
    for run in p.runs:
        run.font.size=Pt(font_size); run.font.bold=bold; run.font.color.rgb=color
    return box

def add_title(slide, title, subtitle=None):
    add_textbox(slide, 0.35, 0.18, 12.3, 0.45, title, 22, True, RGB_NAVY)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.35), Inches(0.72), Inches(12.65), Inches(0.03))
    line.fill.solid(); line.fill.fore_color.rgb = RGB_NAVY; line.line.color.rgb=RGB_NAVY
    if subtitle:
        add_textbox(slide, 0.38, 0.78, 12.2, 0.28, subtitle, 9, False, RGBColor(100,100,100))

def add_footer(slide, source):
    add_textbox(slide, 0.35, 7.12, 12.2, 0.22, f'Source: {source}', 7, False, RGBColor(110,110,110))
    add_textbox(slide, 12.25, 7.12, 0.7, 0.22, str(len(prs.slides)), 7, False, RGBColor(110,110,110), PP_ALIGN.RIGHT)

def add_bullets(slide, x, y, w, h, bullets, font_size=14, color=RGB_DARK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap=True
    for i,b in enumerate(bullets):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(5)
    return box

def add_card(slide, x,y,w,h, headline, value, note='', fill=RGB_LTGREY):
    shape=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb=fill; shape.line.color.rgb=RGBColor(210,210,210)
    add_textbox(slide,x+0.08,y+0.08,w-0.16,0.22,headline,8,True,RGBColor(90,90,90))
    add_textbox(slide,x+0.08,y+0.35,w-0.16,0.40,value,18,True,RGB_NAVY)
    if note: add_textbox(slide,x+0.08,y+0.78,w-0.16,h-0.82,note,7,False,RGBColor(85,85,85))

def add_table(slide, x,y,w,h, data, header=True, font_size=8):
    rows=len(data); cols=len(data[0])
    table = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h)).table
    for r,row in enumerate(data):
        for c,val in enumerate(row):
            cell=table.cell(r,c); cell.text=str(val)
            cell.margin_left=Pt(2); cell.margin_right=Pt(2); cell.margin_top=Pt(1); cell.margin_bottom=Pt(1)
            for p in cell.text_frame.paragraphs:
                p.font.size=Pt(font_size); p.font.color.rgb=RGB_DARK
                if r==0 and header: p.font.bold=True; p.font.color.rgb=RGB_WHITE
            if r==0 and header:
                cell.fill.solid(); cell.fill.fore_color.rgb=RGB_NAVY
            elif r%2==1:
                cell.fill.solid(); cell.fill.fore_color.rgb=RGBColor(248,248,248)
    return table

# Slide 1: Title
s=prs.slides.add_slide(blank)
# dark title band
band=s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
band.fill.solid(); band.fill.fore_color.rgb=RGB_NAVY; band.line.color.rgb=RGB_NAVY
add_textbox(s,0.75,0.85,11.8,0.65,'Shopify Inc. (SHOP)',34,True,RGB_WHITE)
add_textbox(s,0.75,1.55,11.8,0.45,'Investment Committee Presentation',22,False,RGB_WHITE)
add_textbox(s,0.75,2.2,10.8,0.55,'Neutral / Accumulate on pullbacks — exceptional company, valuation still requires proof that 20%+ growth and margin expansion persist',18,False,RGB_WHITE)
add_card(s,0.75,3.2,2.2,1.15,'Reference price',f'${price:.0f}','Koyfin / model',RGBColor(230,238,247))
add_card(s,3.2,3.2,2.2,1.15,'Model DCF value',f'${dcf_value:.0f}','model.xlsx DCF!B107',RGBColor(230,238,247))
add_card(s,5.65,3.2,2.2,1.15,'Implied downside',f'{upside:.1%}','model.xlsx DCF!B109',RGBColor(230,238,247))
add_card(s,8.1,3.2,2.2,1.15,'Market cap / EV','$134B / $128B','Koyfin / model',RGBColor(230,238,247))
add_textbox(s,0.75,6.8,11.5,0.25,'Prepared from SHOP initial coverage report, thesis tracker and DCF model | 2026-05-23',9,False,RGB_WHITE)

# Slide 2: IC recommendation
s=prs.slides.add_slide(blank); add_title(s,'IC recommendation: own the work, wait for either a better entry or proof of reacceleration')
add_card(s,0.6,1.15,2.5,1.1,'Recommendation','Neutral','Accumulate on pullbacks below ~$90 or after proof high-20s growth is a floor',RGBColor(230,238,247))
add_card(s,3.3,1.15,2.5,1.1,'Base DCF','~$96/share','9.5% WACC; 3.5% terminal g; 22% terminal FCF margin',RGBColor(230,238,247))
add_card(s,6.0,1.15,2.5,1.1,'Upside setup','$110–150+','Requires lower WACC, higher growth duration, or AI/B2B upside',RGBColor(230,238,247))
add_card(s,8.7,1.15,2.5,1.1,'Main concern','Valuation','8.3x NTM EV/sales and ~44x NTM EV/EBITDA',RGBColor(230,238,247))
add_bullets(s,0.75,2.75,5.7,3.3,[
    'Business quality is not the debate: Shopify is scaled, founder-led, asset-light and deeply embedded in merchant workflows.',
    'Fundamentals remain strong: Q1 revenue +34%, GMV +35%, FY2025 FCF $2.0B and net cash balance sheet.',
    'But the current multiple still prices in sustained 20%+ revenue growth and FCF margin expansion.',
    'IC action: keep active coverage; add only on price weakness or evidence that growth/margins are tracking above base case.'
],14)
s.shapes.add_picture(str(val_chart), Inches(6.75), Inches(2.55), Inches(5.9), Inches(3.25))
add_footer(s,'report.md; valuation_framework.md; model.xlsx DCF!B107/B109; outputs.json')

# Slide 3: Company snapshot
s=prs.slides.add_slide(blank); add_title(s,'Company snapshot: Shopify is a scaled commerce infrastructure platform, not just SMB store software')
add_card(s,0.6,1.05,2.2,1.05,'FY2025 GMV','$378B','Q1 2026 GMV crossed $100B',RGBColor(242,246,252))
add_card(s,3.0,1.05,2.2,1.05,'FY2025 Revenue','$11.6B','+30% YoY',RGBColor(242,246,252))
add_card(s,5.4,1.05,2.2,1.05,'FY2025 FCF','$2.0B','17% margin',RGBColor(242,246,252))
add_card(s,7.8,1.05,2.2,1.05,'Net cash','~$5.6B','Minimal debt',RGBColor(242,246,252))
add_card(s,10.2,1.05,2.2,1.05,'Insider ownership','~6%','Founder alignment',RGBColor(242,246,252))
s.shapes.add_picture(str(kpi_chart), Inches(0.75), Inches(2.5), Inches(5.75), Inches(3.7))
add_bullets(s,6.85,2.55,5.7,3.5,[
    'Two revenue engines: high-margin Subscription Solutions plus GMV-linked Merchant Solutions.',
    'Merchant Solutions is now ~76% of revenue and is the key monetization flywheel through Payments, Shop Pay, Capital, Shipping, Tax, Markets and POS.',
    'Core strategic pivot: from online storefront software to operating system for independent commerce across online, offline, B2B, international and AI channels.'
],14)
add_footer(s,'Koyfin extraction summary; report.md; thesis_tracker.md')

# Slide 4: Business model flywheel
s=prs.slides.add_slide(blank); add_title(s,'Business model: subscription locks in the merchant; Merchant Solutions monetizes merchant success')
# flow diagram
nodes=[('More merchants',0.9,1.45),('More surfaces\nOnline / POS / B2B / AI',3.4,1.45),('More buyer conversion\nShop Pay / Catalog',6.15,1.45),('More GMV',8.9,1.45),('More Merchant\nSolutions attach',5.3,3.75)]
for text,x,y in nodes:
    shp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.0), Inches(0.85))
    shp.fill.solid(); shp.fill.fore_color.rgb=RGBColor(230,238,247); shp.line.color.rgb=RGB_BLUE
    add_textbox(s,x+0.08,y+0.12,1.85,0.5,text,11,True,RGB_NAVY,PP_ALIGN.CENTER)
# arrows as lines
for x1,y1,x2,y2 in [(2.9,1.88,3.4,1.88),(5.4,1.88,6.15,1.88),(8.15,1.88,8.9,1.88),(9.9,2.3,6.3,3.75),(5.3,4.17,1.9,2.3)]:
    line=s.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2)); line.line.color.rgb=RGB_BLUE; line.line.width=Pt(2)
add_table(s,0.85,5.1,11.6,1.0,[['Revenue stream','Q1 2026 revenue','Gross margin profile','Strategic role'],['Subscription Solutions','$750M / ~24% of revenue','~80%','Merchant acquisition, platform control, Plus/enterprise'],['Merchant Solutions','$2.42B / ~76% of revenue','~39%','GMV-linked monetization and retention flywheel']],True,8)
add_footer(s,'report.md business model section; web_business_strategy.md; Koyfin summary')

# Slide 5: Thesis pillars
s=prs.slides.add_slide(blank); add_title(s,'Thesis pillars: five things must stay true for SHOP to compound through the multiple')
headers=['Pillar','Evidence today','KPI / breakpoint']
data=[headers,
['Commerce OS','FY2025 GMV $378B; Q1 GMV $100.7B; broad app ecosystem','GMV growth; breakpoint <15% for 2+ quarters'],
['Merchant Solutions flywheel','Merchant Solutions 76% of revenue; GPV penetration 67%','GPV penetration; breakpoint stalls <67%'],
['Upmarket / B2B / international','Plus ~34% MRR; B2B +80–96%; international +mid-30s','Plus mix, B2B GMV, enterprise wins'],
['FCF durability','FY2025 FCF $2.0B; Q1 FCF margin 15%','FCF margin; breakpoint <12%'],
['AI option','UCP, Sidekick, Catalog, agentic storefronts','AI-referred GMV / monetization']]
add_table(s,0.55,1.25,12.2,4.7,data,True,8)
add_bullets(s,0.8,6.15,11.7,0.6,['IC framing: thesis is intact today, but valuation requires the KPIs to remain strong through the next 2–4 quarters.'],13)
add_footer(s,'thesis_tracker.md')

# Slide 6: Financial model
s=prs.slides.add_slide(blank); add_title(s,'Financial model: consensus supports 20%+ growth through FY2028; DCF assumes gradual fade and FCF margin expansion')
s.shapes.add_picture(str(rev_chart), Inches(0.65), Inches(1.05), Inches(7.1), Inches(4.0))
add_table(s,8.0,1.05,4.8,2.45,[['Metric','FY2025','FY2028E','FY2035E'],['Revenue','$11.6B','$22.9B','$49.9B'],['Growth','30%','25%','5%'],['FCF margin','17.4%','18.5%','22.0%'],['FCF','$2.0B','$4.2B','$11.0B']],True,9)
add_bullets(s,8.05,3.95,4.55,1.8,[
    'Base case uses Koyfin consensus for FY2026–FY2028.',
    'Long-run value depends on whether Shopify can keep growth elevated while expanding FCF margin despite Merchant Solutions mix pressure.',
    'SBC remains an economic cost; buyback helps but does not eliminate dilution risk.'
],11)
add_footer(s,'model.xlsx DCF tab; valuation_framework.md; report.md financial analysis')

# Slide 7: Valuation
s=prs.slides.add_slide(blank); add_title(s,'Valuation: model base is below current price; upside requires either duration, margin or discount-rate upside')
add_table(s,0.65,1.0,5.9,2.1,[['DCF bridge','Value'],['PV explicit FCF','$39.7B'],['PV terminal value','$80.0B'],['Enterprise value','$119.6B'],['+ Net cash','$5.6B'],['Equity value','$125.2B'],['Implied value / share','$96.30']],True,9)
s.shapes.add_picture(str(heat_chart), Inches(7.05), Inches(0.95), Inches(4.7), Inches(3.15))
add_bullets(s,0.8,3.65,5.6,2.25,[
    f'Model base: ${dcf_value:.2f}/share vs current ${price:.2f}, or {upside:.1%} implied downside.',
    'Base assumptions: 9.5% WACC, 3.5% terminal growth, 22% terminal FCF margin.',
    'Practical decision range: <$85 attractive if fundamentals intact; >$135 requires bull-case execution.'
],12)
add_bullets(s,7.15,4.55,4.8,1.2,[
    'Sensitivity remains large: small changes in WACC / terminal growth can move fair value by >$25/share.',
    'This is why IC should focus on evidence of growth duration and margin trajectory, not just one-point DCF.'
],11)
add_footer(s,'model.xlsx DCF!B100:B109; outputs.json; valuation_framework.md')

# Slide 8: Multiple sanity check
s=prs.slides.add_slide(blank); add_title(s,'Multiple sanity check: SHOP deserves a premium, but the premium is still the key risk')
s.shapes.add_picture(str(mult_chart), Inches(0.65), Inches(1.1), Inches(6.25), Inches(3.75))
add_table(s,7.2,1.05,5.2,2.2,[['Multiple','SHOP current'],['EV/Sales LTM','10.3x'],['EV/Sales NTM','8.3x'],['EV/EBITDA NTM','44.4x'],['P/E NTM','54.5x'],['Street avg target','$151']],True,10)
add_bullets(s,7.25,3.75,5.0,1.9,[
    'SHOP is not directly comparable to small commerce SaaS: scale, growth, GMV, payments and ecosystem justify a premium.',
    'But multiple compression remains the main downside if Q2 guide proves to be the start of a deceleration cycle.',
    'Upside to $130–150 requires sustained premium sales multiple or confidence FY2028 EBITDA/FCF is too low.'
],11)
add_footer(s,'Koyfin multiples; valuation_framework.md comps section; model.xlsx Comps tab')

# Slide 9: Competitive position
s=prs.slides.add_slide(blank); add_title(s,'Competitive position: Shopify wins on merchant alignment, ecosystem breadth and unified commerce')
add_table(s,0.55,1.0,12.2,3.1,[['Competitor set','Where they compete','SHOP counter-positioning'],['Amazon / marketplaces','Demand capture and fulfillment','Shopify preserves merchant brand, customer data and channel independence'],['WooCommerce / Wix / Squarespace','SMB storefront creation','Shopify has deeper commerce stack, payments, apps and scaling path'],['BigCommerce','Mid-market / B2B, no transaction fees','Shopify has larger ecosystem, Shop Pay and broader Merchant Solutions'],['Adobe / Salesforce / custom','Complex enterprise commerce','Shopify offers faster implementation and lower TCO, but must prove high-end complexity'],['Stripe / PayPal / Adyen / Block','Payments / POS economics','Shopify embeds payments in merchant OS; dependence and take-rate risk remain']],True,8)
add_bullets(s,0.75,4.55,11.8,1.3,[
    'Moat source: ecosystem + switching costs + aligned incentives; Shopify does not compete with merchants.',
    'Main competitive risk: AI agents and platform walled gardens could shift demand/checkout away from merchant storefronts.',
    'Monitoring: enterprise wins, GPV penetration, payment take rate, AI checkout economics.'
],13)
add_footer(s,'report.md competitive landscape; web_competition_market.md; thesis_tracker.md')

# Slide 10: Catalysts
s=prs.slides.add_slide(blank); add_title(s,'Catalysts: next 6–12 months should clarify whether deceleration fears are overdone')
add_table(s,0.55,1.0,12.2,4.6,[['Catalyst','Window','Bull signal','Bear signal'],['Q2 2026 earnings','Early Aug 2026','Revenue > guide; GMV resilient; FCF mid/high teens','In-line on lower growth; weak Q3 guide'],['AI / agentic disclosures','2026–2027','AI-referred GMV disclosed; Shopify captures checkout/payment','Traffic without monetization; agents bypass storefront'],['Payments penetration','Quarterly','GPV penetration toward 70–75%','Stalls at 67%; take-rate compression'],['B2B / enterprise','Quarterly','More global brand wins; B2B remains >50% growth','B2B slows; Adobe/Salesforce defend enterprise'],['Buyback execution','2026','Share count stable/down; disciplined price','Buybacks fail to offset SBC dilution'],['Sezzle / regulation','H2 2026+','Case contained','Forced interoperability / payments scrutiny']],True,8)
add_footer(s,'thesis_tracker.md catalysts table')

# Slide 11: Risks and kill criteria
s=prs.slides.add_slide(blank); add_title(s,'Risks and kill criteria: the thesis breaks if growth slows before margins scale')
add_table(s,0.55,1.0,12.2,4.9,[['Risk','Why it matters','Watch item / kill criterion'],['Growth deceleration','Premium multiple requires long growth duration','Revenue <20% or GMV <18% for multiple quarters'],['Merchant Solutions margin pressure','Lower-margin mix can cap FCF expansion','Gross margin <45% or transaction losses rising'],['AI disintermediation','Could bypass storefront / checkout economics','AI traffic grows but Shopify monetization absent'],['Competition','Pressure across SMB, enterprise, POS and payments','Enterprise wins slow; GPV penetration stalls'],['Regulation / antitrust / BNPL','Can impair payments economics','Sezzle case expands or forced payment interoperability'],['SBC dilution','FCF per share matters more than FCF dollars','SBC >5% revenue and buyback insufficient']],True,8)
add_bullets(s,0.8,6.15,11.6,0.55,['Position sizing should reflect asymmetric valuation risk: downside is mostly multiple compression, not balance sheet stress.'],13)
add_footer(s,'thesis_tracker.md risks table; report.md risks section')

# Slide 12: IC decision
s=prs.slides.add_slide(blank); add_title(s,'IC decision: keep on active watch; require either price discipline or KPI confirmation')
add_card(s,0.75,1.15,3.0,1.2,'Decision','Do not initiate full-size position today','Valuation does not provide enough margin of safety at ~$103',RGBColor(242,246,252))
add_card(s,4.05,1.15,3.0,1.2,'Entry rule','Add below ~$90','Assuming GMV / revenue / FCF KPIs remain intact',RGBColor(242,246,252))
add_card(s,7.35,1.15,3.0,1.2,'Upgrade rule','Buy on proof','Q2/Q3 show high-20s growth floor + FCF durability',RGBColor(242,246,252))
add_bullets(s,0.95,3.05,5.6,2.4,[
    'What we like: founder-led, dominant independent commerce OS, GMV-linked monetization, net cash, FCF-positive, multiple growth vectors.',
    'What holds us back: base DCF below current price, premium multiple, margin mix pressure, AI outcome uncertain.',
    'What to do now: maintain coverage, update after Q2 earnings, track KPI dashboard and valuation bands.'
],13)
s.shapes.add_picture(str(val_chart), Inches(6.9), Inches(3.0), Inches(5.55), Inches(2.8))
add_footer(s,'Investment conclusion from report.md; model.xlsx; thesis_tracker.md')

prs.save(OUT)
print(json.dumps({'status':'success','deck':str(OUT),'slides':len(prs.slides),'charts':[str(p) for p in [rev_chart,val_chart,kpi_chart,mult_chart,heat_chart]]}, indent=2))
